"""Build a leave-behind information deck for Standard Motor Co.

A self-contained reference deck the recipient can read on their own
after the conversation. Same architectural content as the live-meeting
briefing, but with founder-internal language stripped out and a closing
resources slide added so the recipient can verify claims independently.

Run: python3 build_standard_motor_info_deck.py
Output: ../info-deck-standard-motor-2026-04-28.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Theme ────────────────────────────────────────────────────────────────

NAVY = RGBColor(0x0B, 0x19, 0x33)
INK = RGBColor(0x1A, 0x1F, 0x36)
GREY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
LINE = RGBColor(0xE2, 0xE5, 0xEB)

ACCENT_RED = RGBColor(0xE5, 0x4B, 0x4B)
ACCENT_AMBER = RGBColor(0xE0, 0x9F, 0x3E)
ACCENT_PURPLE = RGBColor(0x8E, 0x6F, 0xE6)
ACCENT_BLUE = RGBColor(0x4F, 0x8F, 0xE5)
ACCENT_GREEN = RGBColor(0x36, 0xA1, 0x6B)

FONT = "Calibri"

DECK_TITLE = "AI Identity"
DECK_SUBTITLE = "Standard Motor Co briefing"
DECK_FOOTER = "AI Identity  •  Information deck  •  Standard Motor Co"
TODAY = "April 28, 2026"
FOUNDER_EMAIL = "jeff@ai-identity.co"
FOUNDER_NAME = "Jeff Leva"
WEBSITE = "ai-identity.co"

OUT = Path(__file__).resolve().parent.parent / "info-deck-standard-motor-2026-04-28.pptx"

# ── Helpers ──────────────────────────────────────────────────────────────


def set_text(tf, runs, *, anchor=MSO_ANCHOR.TOP):
    """Replace a text frame's content with a list of runs.
    runs = [(text, {size, bold, color, font, align}), ...]
    The first run becomes the first paragraph; subsequent runs each
    become a new paragraph unless align='inline'.
    """
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    for i, (text, opts) in enumerate(runs):
        if opts.get("align") == "inline" and i > 0:
            r = p.add_run()
        else:
            if i > 0:
                p = tf.add_paragraph()
            r = p.runs[0] if p.runs else p.add_run()
        if opts.get("align") in ("center", "left", "right"):
            p.alignment = {
                "center": PP_ALIGN.CENTER,
                "left": PP_ALIGN.LEFT,
                "right": PP_ALIGN.RIGHT,
            }[opts["align"]]
        r.text = text
        r.font.name = opts.get("font", FONT)
        r.font.size = Pt(opts.get("size", 14))
        r.font.bold = opts.get("bold", False)
        r.font.color.rgb = opts.get("color", INK)
        if "space_after" in opts:
            p.space_after = Pt(opts["space_after"])


def add_rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    return s


def add_text(slide, x, y, w, h, runs, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    set_text(box.text_frame, runs, anchor=anchor)
    return box


def add_card(slide, x, y, w, h, *, accent, title, body, body_size=12):
    """Light grey rounded card with a colored left bar."""
    add_rect(slide, x, y, w, h, fill=LIGHT_BG, line=LINE)
    add_rect(slide, x, y, Emu(60000), h, fill=accent)
    add_text(
        slide,
        x + Inches(0.25),
        y + Inches(0.2),
        w - Inches(0.4),
        Inches(0.5),
        [(title, dict(size=15, bold=True, color=NAVY))],
    )
    add_text(
        slide,
        x + Inches(0.25),
        y + Inches(0.75),
        w - Inches(0.4),
        h - Inches(0.85),
        [(body, dict(size=body_size, color=INK))],
    )


def add_footer(slide, page, total):
    add_text(
        slide,
        Inches(0.5),
        Inches(7.05),
        Inches(12.33),
        Inches(0.3),
        [(f"{DECK_FOOTER}  •  {page} / {total}", dict(size=9, color=GREY, align="right"))],
    )


def add_slide_title(slide, title, subtitle=None):
    add_text(
        slide,
        Inches(0.6),
        Inches(0.45),
        Inches(12.13),
        Inches(0.9),
        [(title, dict(size=34, bold=True, color=NAVY))],
    )
    if subtitle:
        add_text(
            slide,
            Inches(0.6),
            Inches(1.25),
            Inches(12.13),
            Inches(0.55),
            [(subtitle, dict(size=15, color=GREY))],
        )
    add_rect(slide, Inches(0.6), Inches(1.85), Inches(12.13), Pt(1), fill=LINE)


# ── Build ────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]
SLIDES_PLANNED = 16

# 1. Title ----------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, fill=NAVY)
add_rect(s, Inches(0.6), Inches(2.4), Inches(0.18), Inches(2.5), fill=ACCENT_AMBER)
add_text(
    s,
    Inches(1.0),
    Inches(2.3),
    Inches(11.5),
    Inches(1.0),
    [("AI Identity", dict(size=68, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF)))],
)
add_text(
    s,
    Inches(1.0),
    Inches(3.5),
    Inches(11.5),
    Inches(0.5),
    [
        (
            "Cryptographic identity, policy, and audit for AI agents",
            dict(size=22, color=RGBColor(0xC9, 0xD4, 0xE6)),
        )
    ],
)
add_text(
    s,
    Inches(1.0),
    Inches(4.7),
    Inches(11.5),
    Inches(0.5),
    [
        (
            f"Prepared for Kevin Pentecost  ·  Standard Motor Co  ·  {TODAY}",
            dict(size=14, color=RGBColor(0x9D, 0xA9, 0xC1)),
        )
    ],
)
add_text(
    s,
    Inches(1.0),
    Inches(6.55),
    Inches(11.5),
    Inches(0.4),
    [
        (
            f"{FOUNDER_NAME}  ·  Founder & CEO  ·  {FOUNDER_EMAIL}  ·  {WEBSITE}",
            dict(size=11, color=RGBColor(0xC9, 0xD4, 0xE6), align="left"),
        )
    ],
)

# 2. The 30-second pitch --------------------------------------------------
s = prs.slides.add_slide(blank)
add_slide_title(s, "What AI Identity is", "The one-sentence version")
add_rect(s, Inches(0.6), Inches(2.3), Inches(0.12), Inches(2.6), fill=ACCENT_AMBER)
add_text(
    s,
    Inches(1.0),
    Inches(2.4),
    Inches(11.7),
    Inches(2.5),
    [
        (
            "“AI Identity is the cryptographic identity, policy, and audit layer for AI agents — "
            "built from scratch for autonomous systems, not retrofitted from human IAM.”",
            dict(size=26, color=NAVY, bold=True),
        )
    ],
)
add_text(
    s,
    Inches(1.0),
    Inches(5.1),
    Inches(11.7),
    Inches(0.4),
    [("What's inside this deck:", dict(size=14, color=GREY, bold=True))],
)
for i, line in enumerate(
    [
        "1.  The problem every security team is about to inherit",
        "2.  How teams cope today, and why those workarounds don't actually work",
        "3.  How AI Identity solves it — gateway, three-key architecture, forensic chain",
        "4.  Two illustrative scenarios: parts ordering, customer service refunds",
        "5.  Common questions, and where to verify the technical claims yourself",
    ]
):
    add_text(
        s,
        Inches(1.0),
        Inches(5.5) + Inches(0.27) * i,
        Inches(11.7),
        Inches(0.3),
        [(line, dict(size=12, color=INK))],
    )
add_footer(s, 2, SLIDES_PLANNED)

# 3. The Problem (original framing, sharper) ------------------------------
s = prs.slides.add_slide(blank)
add_slide_title(
    s,
    "The Problem",
    "AI agents are in production now — reading databases, calling APIs, moving money. Most companies treat them like anonymous scripts.",
)
card_w = Inches(5.95)
card_h = Inches(2.3)
left_x = Inches(0.6)
right_x = Inches(6.78)
top_y = Inches(2.15)
bot_y = Inches(4.65)
add_card(
    s,
    left_x,
    top_y,
    card_w,
    card_h,
    accent=ACCENT_RED,
    title="No Identity",
    body="You can't tell which agent took which action. Every agent shares the same service-account key — the audit log records the key, not the persona behind it.",
)
add_card(
    s,
    right_x,
    top_y,
    card_w,
    card_h,
    accent=ACCENT_AMBER,
    title="No Policy Enforcement",
    body="Agents get over-permissioned during the demo, and nothing gets removed. Permission drift becomes the steady state.",
)
add_card(
    s,
    left_x,
    bot_y,
    card_w,
    card_h,
    accent=ACCENT_PURPLE,
    title="No Forensic Record",
    body="When something goes wrong, there's no tamper-proof log to show a regulator. Standard application logs are mutable — anyone with write access can edit them.",
)
add_card(
    s,
    right_x,
    bot_y,
    card_w,
    card_h,
    accent=ACCENT_BLUE,
    title="Traditional IAM Doesn't Fit",
    body="Okta and Azure Entra are bolting agent features onto human-first architectures. They don't capture chain-of-thought, can't replay sessions, and treat agents as service accounts.",
)
add_footer(s, 3, SLIDES_PLANNED)

# 4. What teams do today (workaround table) -------------------------------
s = prs.slides.add_slide(blank)
add_slide_title(
    s,
    "What teams do today",
    "Every workaround has the same shape: the thing being audited is also the thing asserting its identity.",
)
rows = [
    (
        "Self-asserted metadata field",
        "Pass user_id / agent_name on every API call. Filter logs by it.",
        "The agent declares its own name. A compromised agent can claim to be anyone.",
        ACCENT_RED,
    ),
    (
        "One API key per agent",
        "Generate a separate Anthropic / OpenAI key for each agent. Keys act as identity.",
        "Keys never expire, no central revocation, secret sprawl. Every rotation is a multi-system deploy.",
        ACCENT_RED,
    ),
    (
        "LLM observability proxies",
        "LangSmith, Langfuse, Helicone, Arize, Datadog LLM Obs. Tag each request with userId / sessionId.",
        "Same self-assertion — tags come from the calling code. Now you trust the proxy too.",
        ACCENT_AMBER,
    ),
    (
        "Workspaces / Projects",
        "Anthropic Workspaces, OpenAI Projects. Coarse-grained per-team partitioning.",
        "Useful for billing. Doesn't answer 'which agent inside the Workspace did this.'",
        ACCENT_AMBER,
    ),
    (
        "Internal gateway (homegrown)",
        "Route every model call through an in-house proxy. Verify mTLS or internal JWT.",
        "Verifies the workload, not the agent. Two agents in the same pod look identical. Quarter of platform work.",
        ACCENT_BLUE,
    ),
]
y = Inches(2.05)
row_h = Inches(0.92)
gutter = Inches(0.04)
header_color = GREY
add_text(
    s,
    Inches(0.6),
    Inches(2.0),
    Inches(3.5),
    Inches(0.3),
    [("PATTERN", dict(size=10, bold=True, color=header_color))],
)
add_text(
    s,
    Inches(4.2),
    Inches(2.0),
    Inches(4.5),
    Inches(0.3),
    [("HOW IT WORKS", dict(size=10, bold=True, color=header_color))],
)
add_text(
    s,
    Inches(8.85),
    Inches(2.0),
    Inches(3.9),
    Inches(0.3),
    [("WHY IT FALLS SHORT", dict(size=10, bold=True, color=header_color))],
)
y = Inches(2.35)
for pattern, how, gap, color in rows:
    add_rect(s, Inches(0.6), y, Inches(12.13), row_h, fill=LIGHT_BG, line=LINE)
    add_rect(s, Inches(0.6), y, Emu(50000), row_h, fill=color)
    add_text(
        s,
        Inches(0.78),
        y + Inches(0.1),
        Inches(3.4),
        row_h - Inches(0.2),
        [(pattern, dict(size=12, bold=True, color=NAVY))],
    )
    add_text(
        s,
        Inches(4.2),
        y + Inches(0.1),
        Inches(4.55),
        row_h - Inches(0.2),
        [(how, dict(size=10, color=INK))],
    )
    add_text(
        s,
        Inches(8.85),
        y + Inches(0.1),
        Inches(3.78),
        row_h - Inches(0.2),
        [(gap, dict(size=10, color=INK))],
    )
    y += row_h + gutter
add_footer(s, 4, SLIDES_PLANNED)

# 5. The fundamental gap --------------------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, fill=NAVY)
add_text(
    s,
    Inches(1.0),
    Inches(1.6),
    Inches(11.3),
    Inches(0.5),
    [("THE FUNDAMENTAL GAP", dict(size=14, bold=True, color=ACCENT_AMBER, align="left"))],
)
add_text(
    s,
    Inches(1.0),
    Inches(2.3),
    Inches(11.3),
    Inches(2.6),
    [
        (
            "“The thing being audited is also the thing asserting its identity. "
            "No external verification.”",
            dict(size=34, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF)),
        )
    ],
)
add_text(
    s,
    Inches(1.0),
    Inches(5.2),
    Inches(11.3),
    Inches(1.4),
    [
        (
            "It's the equivalent of letting every employee write whatever name they want on the office sign-in sheet, "
            "and then pulling the sheet to investigate a break-in. AI Identity flips this — identity is verified by an "
            "external party (the AI Identity gateway), bound to a hardware-rooted key the agent never holds.",
            dict(size=15, color=RGBColor(0xC9, 0xD4, 0xE6)),
        )
    ],
)
add_text(
    s,
    Inches(0.6),
    Inches(7.05),
    Inches(12.13),
    Inches(0.3),
    [
        (
            f"{DECK_FOOTER}  •  5 / {SLIDES_PLANNED}",
            dict(size=9, color=RGBColor(0x6F, 0x7D, 0x9C), align="right"),
        )
    ],
)

# 6. Four Pillars ---------------------------------------------------------
s = prs.slides.add_slide(blank)
add_slide_title(
    s,
    "What AI Identity does",
    "Four pillars — every enterprise must answer four questions about its AI agents.",
)
pillars = [
    (
        "Identity",
        "Who is this agent?",
        "Per-agent aid_sk_ keys. Issued, rotated, revoked, expired. SHA-256 hashed at rest. Compromise one — kill one.",
        ACCENT_RED,
    ),
    (
        "Policy",
        "What is it allowed to do?",
        "Deny-by-default gateway. Per-agent scopes, spend caps, allowed endpoints, allowed data. Enforced at runtime, not in a config file.",
        ACCENT_AMBER,
    ),
    (
        "Forensics",
        "What did it actually do, and why?",
        "HMAC-chained audit log. KMS-signed evidence packets. Step through any session — auth, policy checks, tool calls, blocks.",
        ACCENT_PURPLE,
    ),
    (
        "Compliance",
        "Can you prove it to a regulator?",
        "SOC 2 Type II, EU AI Act, NIST AI RMF, GDPR profiles. One-click export. Continuous, not annual.",
        ACCENT_BLUE,
    ),
]
card_w = Inches(2.95)
card_h = Inches(4.5)
x = Inches(0.6)
gap = Inches(0.12)
for title, sub, body, color in pillars:
    add_rect(s, x, Inches(2.15), card_w, card_h, fill=LIGHT_BG, line=LINE)
    add_rect(s, x, Inches(2.15), card_w, Inches(0.12), fill=color)
    add_text(
        s,
        x + Inches(0.25),
        Inches(2.4),
        card_w - Inches(0.4),
        Inches(0.55),
        [(title, dict(size=22, bold=True, color=NAVY))],
    )
    add_text(
        s,
        x + Inches(0.25),
        Inches(3.0),
        card_w - Inches(0.4),
        Inches(0.5),
        [(sub, dict(size=12, color=color, bold=True))],
    )
    add_text(
        s,
        x + Inches(0.25),
        Inches(3.55),
        card_w - Inches(0.4),
        card_h - Inches(1.6),
        [(body, dict(size=11, color=INK))],
    )
    x += card_w + gap
add_footer(s, 6, SLIDES_PLANNED)

# 7. The gateway: <50ms enforcement ---------------------------------------
s = prs.slides.add_slide(blank)
add_slide_title(
    s,
    "The gateway: <50ms enforcement",
    "Every agent request runs through a six-stage fail-closed pipeline. Decision target: under 50 milliseconds.",
)
stages = [
    ("Rate Limit", "100/s per IP\n60/s per agent", "429"),
    ("Key Validation", "aid_sk_ verified\nagainst KMS-hashed registry", "401"),
    ("Agent Status", "active / suspended /\nrevoked / expired", "403"),
    ("Policy Match", "deny-by-default\ncontext-aware", "403"),
    ("Circuit Breaker", "5 failures / 60s →\nOPEN until probe", "503"),
    ("Audit Log", "HMAC-chained,\nappend-only", "200"),
]
stage_w = Inches(1.95)
stage_h = Inches(2.0)
x = Inches(0.6)
y = Inches(2.4)
for i, (name, body, code) in enumerate(stages):
    color = ACCENT_GREEN if i == len(stages) - 1 else ACCENT_BLUE
    add_rect(s, x, y, stage_w, stage_h, fill=LIGHT_BG, line=LINE)
    add_rect(s, x, y, stage_w, Inches(0.12), fill=color)
    add_text(
        s,
        x + Inches(0.18),
        y + Inches(0.25),
        stage_w - Inches(0.3),
        Inches(0.55),
        [(name, dict(size=13, bold=True, color=NAVY))],
    )
    add_text(
        s,
        x + Inches(0.18),
        y + Inches(0.85),
        stage_w - Inches(0.3),
        Inches(0.85),
        [(body, dict(size=9.5, color=INK))],
    )
    add_text(
        s,
        x + Inches(0.18),
        y + stage_h - Inches(0.45),
        stage_w - Inches(0.3),
        Inches(0.4),
        [(f"on fail → {code}", dict(size=8.5, color=GREY))],
    )
    if i < len(stages) - 1:
        arrow = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x + stage_w + Inches(0.005),
            y + Inches(0.85),
            Inches(0.2),
            Inches(0.3),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = LINE
        arrow.line.fill.background()
    x += stage_w + Inches(0.205)
add_text(
    s,
    Inches(0.6),
    Inches(4.7),
    Inches(12.13),
    Inches(0.4),
    [("Why <50ms matters operationally:", dict(size=14, bold=True, color=NAVY))],
)
for i, line in enumerate(
    [
        "•  Inline enforcement on the agent's hot path — your agents stay fast.",
        "•  Fail-closed: any exception or timeout = DENY. Security default, not a config switch.",
        "•  4 thread-pool workers, 500ms hard timeout, circuit breaker isolates a misbehaving downstream.",
        "•  Public deployment runs on GKE Autopilot in us-east1 — multi-region path on the H2 roadmap.",
    ]
):
    add_text(
        s,
        Inches(0.8),
        Inches(5.15) + Inches(0.32) * i,
        Inches(12.0),
        Inches(0.3),
        [(line, dict(size=12, color=INK))],
    )
add_footer(s, 7, SLIDES_PLANNED)

# 8. The dual-key (actually triple-key) story -----------------------------
s = prs.slides.add_slide(blank)
add_slide_title(
    s,
    "The dual-key story (actually three keys)",
    "Forging an audit trail would require compromising two separately-held secrets. The hardest one is in an HSM that even we can't read.",
)
keys = [
    (
        "Runtime key  ·  aid_sk_…",
        "Per-agent",
        "Held by the agent. Used to authenticate requests at the gateway. Rotatable. Revocable. SHA-256 hashed at rest. Shown once at creation, never again.",
        "Compromise blast radius: ONE agent.",
        ACCENT_RED,
    ),
    (
        "HMAC chain key  ·  audit_hmac_key",
        "Per-org, in app DB",
        "Signs every audit log row. Each row is HMAC-linked to the previous one — tamper with any record and the chain breaks. Lives in the AI Identity application database.",
        "Compromise blast radius: chain integrity for that org.",
        ACCENT_AMBER,
    ),
    (
        "Forensic signing key  ·  EC_SIGN_P256_SHA256",
        "Per-org, in GCP KMS HSM",
        "Signs every evidence packet at session close. The private half lives in Google's FIPS 140-2 Level 3 hardware. AI Identity itself has never seen it. Public key published via JWKS for offline verification.",
        "Compromise blast radius: requires breaking an HSM.",
        ACCENT_PURPLE,
    ),
]
y = Inches(2.1)
row_h = Inches(1.55)
for title, scope, body, blast, color in keys:
    add_rect(s, Inches(0.6), y, Inches(12.13), row_h, fill=LIGHT_BG, line=LINE)
    add_rect(s, Inches(0.6), y, Emu(60000), row_h, fill=color)
    add_text(
        s,
        Inches(0.85),
        y + Inches(0.18),
        Inches(4.4),
        Inches(0.5),
        [(title, dict(size=14, bold=True, color=NAVY))],
    )
    add_text(
        s,
        Inches(0.85),
        y + Inches(0.62),
        Inches(4.4),
        Inches(0.4),
        [(scope, dict(size=11, color=color, bold=True))],
    )
    add_text(
        s,
        Inches(5.4),
        y + Inches(0.18),
        Inches(7.1),
        Inches(0.95),
        [(body, dict(size=11, color=INK))],
    )
    add_text(
        s,
        Inches(5.4),
        y + Inches(1.15),
        Inches(7.1),
        Inches(0.3),
        [(blast, dict(size=10.5, color=GREY, bold=True))],
    )
    y += row_h + Inches(0.08)
add_text(
    s,
    Inches(0.6),
    Inches(6.95),
    Inches(12.13),
    Inches(0.3),
    [
        (
            "To forge a session, an attacker needs BOTH the HMAC key AND the KMS private key. The KMS private key never leaves the HSM.",
            dict(size=11.5, color=NAVY, bold=True),
        )
    ],
)
add_footer(s, 8, SLIDES_PLANNED)

# 9. What an attestation proves (and doesn't) -----------------------------
s = prs.slides.add_slide(blank)
add_slide_title(
    s,
    "What a signed attestation proves",
    "What you can hand a SOC 2 auditor or a regulator — and what you still need other evidence for.",
)
proves = [
    "The envelope was produced by the key identified in signer_key_id.",
    "At signed_at, AI Identity had those specific audit rows in the state they appear now.",
    "The rows form a contiguous HMAC chain — no insertions, deletions, or modifications since.",
    "Re-ordering or swapping rows is detectable via the chain's tail hash.",
    "The committed range is self-consistent: row count and IDs match.",
]
not_proves = [
    "Whether the agent's *intent* was good — attestation binds rows, rows record gateway decisions.",
    "Whether a denial was *correct by policy* — attestation says what the gateway decided, not whether the policy was right.",
    "Sub-second timing within the session — only signed_at is the signer's timestamp.",
    "Events outside the committed range — needs a separate attestation.",
    "That AI Identity itself wasn't compromised — that's why the HMAC chain uses a *different* key (see prior slide).",
]
add_rect(s, Inches(0.6), Inches(2.15), Inches(5.95), Inches(4.5), fill=LIGHT_BG, line=LINE)
add_rect(s, Inches(0.6), Inches(2.15), Inches(5.95), Inches(0.12), fill=ACCENT_GREEN)
add_text(
    s,
    Inches(0.85),
    Inches(2.4),
    Inches(5.5),
    Inches(0.45),
    [("PROVES", dict(size=14, bold=True, color=ACCENT_GREEN))],
)
for i, line in enumerate(proves):
    add_text(
        s,
        Inches(0.85),
        Inches(2.95) + Inches(0.65) * i,
        Inches(5.5),
        Inches(0.65),
        [(f"✓  {line}", dict(size=11, color=INK))],
    )
add_rect(s, Inches(6.78), Inches(2.15), Inches(5.95), Inches(4.5), fill=LIGHT_BG, line=LINE)
add_rect(s, Inches(6.78), Inches(2.15), Inches(5.95), Inches(0.12), fill=ACCENT_RED)
add_text(
    s,
    Inches(7.03),
    Inches(2.4),
    Inches(5.5),
    Inches(0.45),
    [("DOES NOT PROVE", dict(size=14, bold=True, color=ACCENT_RED))],
)
for i, line in enumerate(not_proves):
    add_text(
        s,
        Inches(7.03),
        Inches(2.95) + Inches(0.65) * i,
        Inches(5.5),
        Inches(0.65),
        [(f"×  {line}", dict(size=11, color=INK))],
    )
add_text(
    s,
    Inches(0.6),
    Inches(6.85),
    Inches(12.13),
    Inches(0.4),
    [
        (
            "Trust model documented in full at ai-identity.co/docs/forensics/trust-model — designed to be auditor-defensible.",
            dict(size=11.5, color=GREY, bold=True),
        )
    ],
)
add_footer(s, 9, SLIDES_PLANNED)

# 10. Standard Motor scenario #1: Parts ordering -------------------------
s = prs.slides.add_slide(blank)
add_slide_title(
    s,
    "Scenario 1: parts ordering AI agent",
    "The kind of system Standard Motor's parts/procurement team is one product cycle away from running in production.",
)
add_text(
    s,
    Inches(0.6),
    Inches(2.05),
    Inches(12.13),
    Inches(0.4),
    [("THE SETUP", dict(size=11, bold=True, color=GREY))],
)
add_text(
    s,
    Inches(0.6),
    Inches(2.4),
    Inches(12.13),
    Inches(0.6),
    [
        (
            "An AI agent monitors inventory, reorders parts when stock falls below threshold, splits orders across approved suppliers, "
            "and updates the ERP. Humans are looped in on orders > $X. The agent has access to supplier APIs, the ERP, and the company AmEx.",
            dict(size=12, color=INK),
        )
    ],
)
risks = [
    (
        "Prompt injection in a supplier email reorders 10× normal qty before anyone notices.",
        ACCENT_RED,
    ),
    ("Compromised key wires payment to a fraudulent supplier added by the attacker.", ACCENT_RED),
    (
        "Over-permissioned agent silently raises its own credit limit during a 'just one more thing' iteration.",
        ACCENT_AMBER,
    ),
    (
        "Auditor asks 'why did this $250K order go through' — no clean evidence trail to answer.",
        ACCENT_AMBER,
    ),
]
fixes = [
    "Per-agent spend cap. The policy denies any order > $X without human-in-the-loop.",
    "Supplier allowlist enforced at the gateway. New supplier ≠ in allowlist = 403, audited.",
    "Per-agent scopes. Credit-limit endpoints are explicitly denied; 'just one more thing' becomes a deploy event, not a runtime change.",
    "KMS-signed evidence packet for every session. Replay the agent's reasoning chain for the auditor in 30 seconds.",
]
add_text(
    s,
    Inches(0.6),
    Inches(3.45),
    Inches(5.95),
    Inches(0.4),
    [("RISKS without us", dict(size=11, bold=True, color=ACCENT_RED))],
)
for i, (text, color) in enumerate(risks):
    add_rect(
        s,
        Inches(0.6),
        Inches(3.85) + Inches(0.7) * i,
        Inches(5.95),
        Inches(0.62),
        fill=LIGHT_BG,
        line=LINE,
    )
    add_rect(s, Inches(0.6), Inches(3.85) + Inches(0.7) * i, Emu(50000), Inches(0.62), fill=color)
    add_text(
        s,
        Inches(0.83),
        Inches(3.85) + Inches(0.7) * i + Inches(0.08),
        Inches(5.6),
        Inches(0.5),
        [(text, dict(size=10.5, color=INK))],
    )
add_text(
    s,
    Inches(6.78),
    Inches(3.45),
    Inches(5.95),
    Inches(0.4),
    [("HOW AI IDENTITY HELPS", dict(size=11, bold=True, color=ACCENT_GREEN))],
)
for i, text in enumerate(fixes):
    add_rect(
        s,
        Inches(6.78),
        Inches(3.85) + Inches(0.7) * i,
        Inches(5.95),
        Inches(0.62),
        fill=LIGHT_BG,
        line=LINE,
    )
    add_rect(
        s, Inches(6.78), Inches(3.85) + Inches(0.7) * i, Emu(50000), Inches(0.62), fill=ACCENT_GREEN
    )
    add_text(
        s,
        Inches(7.01),
        Inches(3.85) + Inches(0.7) * i + Inches(0.08),
        Inches(5.6),
        Inches(0.5),
        [(text, dict(size=10.5, color=INK))],
    )
add_footer(s, 10, SLIDES_PLANNED)

# 11. Standard Motor scenario #2: Customer service refunds ---------------
s = prs.slides.add_slide(blank)
add_slide_title(
    s,
    "Scenario 2: customer service refund agent",
    "A direct-to-customer use case Standard Motor's CX team will face the moment they pilot a tier-1 deflection bot.",
)
add_text(
    s,
    Inches(0.6),
    Inches(2.05),
    Inches(12.13),
    Inches(0.4),
    [("THE SETUP", dict(size=11, bold=True, color=GREY))],
)
add_text(
    s,
    Inches(0.6),
    Inches(2.4),
    Inches(12.13),
    Inches(0.6),
    [
        (
            "An AI agent reviews customer accounts, looks up order history, can issue refunds up to $X without escalation. "
            "Reads from your CRM, has scoped write access to your billing system, and can email customers.",
            dict(size=12, color=INK),
        )
    ],
)
risks = [
    (
        "Prompt injection in a customer message convinces the agent to issue a $1,200 refund on a $50 order.",
        ACCENT_RED,
    ),
    (
        "Agent's session leaks PII — full payment cards, addresses — into the model provider's logs.",
        ACCENT_RED,
    ),
    (
        "Months later, a class action alleges biased refund decisions. Your defense is application logs.",
        ACCENT_AMBER,
    ),
    (
        "Stolen agent credential issues hundreds of small refunds to attacker-controlled accounts overnight.",
        ACCENT_AMBER,
    ),
]
fixes = [
    "Per-agent refund cap enforced at the gateway. Anything above $X is denied or routed to a human.",
    "PII-scoped policies. Card numbers are tokenized before they reach the model provider; outbound calls deny by default.",
    "Cryptographically chained reasoning log. Replay every refund decision; export the SOC 2 evidence package.",
    "Anomaly detection: spike of refund requests from a single agent triggers automatic suspension within seconds.",
]
add_text(
    s,
    Inches(0.6),
    Inches(3.45),
    Inches(5.95),
    Inches(0.4),
    [("RISKS without us", dict(size=11, bold=True, color=ACCENT_RED))],
)
for i, (text, color) in enumerate(risks):
    add_rect(
        s,
        Inches(0.6),
        Inches(3.85) + Inches(0.7) * i,
        Inches(5.95),
        Inches(0.62),
        fill=LIGHT_BG,
        line=LINE,
    )
    add_rect(s, Inches(0.6), Inches(3.85) + Inches(0.7) * i, Emu(50000), Inches(0.62), fill=color)
    add_text(
        s,
        Inches(0.83),
        Inches(3.85) + Inches(0.7) * i + Inches(0.08),
        Inches(5.6),
        Inches(0.5),
        [(text, dict(size=10.5, color=INK))],
    )
add_text(
    s,
    Inches(6.78),
    Inches(3.45),
    Inches(5.95),
    Inches(0.4),
    [("HOW AI IDENTITY HELPS", dict(size=11, bold=True, color=ACCENT_GREEN))],
)
for i, text in enumerate(fixes):
    add_rect(
        s,
        Inches(6.78),
        Inches(3.85) + Inches(0.7) * i,
        Inches(5.95),
        Inches(0.62),
        fill=LIGHT_BG,
        line=LINE,
    )
    add_rect(
        s, Inches(6.78), Inches(3.85) + Inches(0.7) * i, Emu(50000), Inches(0.62), fill=ACCENT_GREEN
    )
    add_text(
        s,
        Inches(7.01),
        Inches(3.85) + Inches(0.7) * i + Inches(0.08),
        Inches(5.6),
        Inches(0.5),
        [(text, dict(size=10.5, color=INK))],
    )
add_footer(s, 11, SLIDES_PLANNED)

# 12. A day in the life ---------------------------------------------------
s = prs.slides.add_slide(blank)
add_slide_title(
    s,
    "What this looks like operationally",
    "How your security team would actually use AI Identity day-to-day.",
)
steps = [
    (
        "Day 1",
        "Engineer registers a new agent.",
        "Single API call → returns aid_sk_ once. Permissions default to deny. Engineer writes a policy: 'this agent can call /reorder, max $5K/day, supplier in {approved-list}.' Reviewed and approved by your security team in the dashboard.",
    ),
    (
        "Day 1, 30 mins later",
        "First request hits the gateway.",
        "Sub-50ms decision: rate limit OK, key valid, agent active, policy match, no circuit breakers tripped. Allowed. Logged. Customer sees a normal API response.",
    ),
    (
        "Day 14",
        "An anomaly fires.",
        "Agent suddenly tries 50 reorders in 60 seconds. Anomaly detection flags it; auto-revoke kicks in; on-call gets paged with the agent ID and the last 100 audit rows pre-filtered.",
    ),
    (
        "Quarterly",
        "SOC 2 audit shows up.",
        "Click 'Export → SOC 2 Type II.' You hand the auditor an attestation envelope, the audit log range it commits to, and the JWKS public key. They verify offline with the open-source CLI. No back-and-forth.",
    ),
]
y = Inches(2.15)
for label, headline, body in steps:
    add_rect(s, Inches(0.6), y, Inches(12.13), Inches(1.1), fill=LIGHT_BG, line=LINE)
    add_rect(s, Inches(0.6), y, Inches(1.7), Inches(1.1), fill=NAVY)
    add_text(
        s,
        Inches(0.7),
        y + Inches(0.35),
        Inches(1.55),
        Inches(0.4),
        [(label, dict(size=12, bold=True, color=ACCENT_AMBER))],
    )
    add_text(
        s,
        Inches(2.45),
        y + Inches(0.13),
        Inches(10.2),
        Inches(0.45),
        [(headline, dict(size=14, bold=True, color=NAVY))],
    )
    add_text(
        s,
        Inches(2.45),
        y + Inches(0.55),
        Inches(10.2),
        Inches(0.55),
        [(body, dict(size=10.5, color=INK))],
    )
    y += Inches(1.18)
add_footer(s, 12, SLIDES_PLANNED)

# 13. How we differ from Okta/Entra/CyberArk ------------------------------
s = prs.slides.add_slide(blank)
add_slide_title(
    s,
    "How we differ from Okta, Entra, CyberArk",
    "They're retrofitting agent features onto human-first architectures. We were built from scratch for autonomous systems.",
)
headers = ["", "Human IAM (Okta / Entra / CyberArk)", "AI Identity"]
rows = [
    (
        "Identity model",
        "User accounts, service accounts, OAuth tokens — built for sessions a human starts.",
        "Per-agent runtime keys, attached to a verified persona, rotated/revoked by API.",
    ),
    (
        "Policy granularity",
        "Role-based: this group can access this app.",
        "Per-decision: this agent, this endpoint, this argument, this spend cap, evaluated at runtime.",
    ),
    (
        "Audit log",
        "Mutable application logs. Tamper-evident only if you bolt on a SIEM with WORM storage.",
        "HMAC-chained at write time. Tamper detection is a property of the log, not of an external system.",
    ),
    (
        "Forensic export",
        "You write a query. The auditor takes your word for it.",
        "DSSE-signed evidence packet + JWKS. The auditor verifies offline against an HSM-backed key.",
    ),
    (
        "Designed for",
        "Humans (millions of them) authenticating to apps.",
        "Agents (millions of them) running with delegated authority and chain-of-thought reasoning.",
    ),
]
y = Inches(2.05)
add_rect(s, Inches(0.6), y, Inches(12.13), Inches(0.4), fill=NAVY)
for i, h in enumerate(headers):
    x = Inches(0.6) + Inches(2.7) * (0 if i == 0 else 1) + (Inches(4.71) if i == 2 else Inches(0))
    if i == 0:
        x = Inches(0.6)
        w = Inches(2.7)
    elif i == 1:
        x = Inches(3.3)
        w = Inches(4.71)
    else:
        x = Inches(8.02)
        w = Inches(4.71)
    add_text(
        s,
        x + Inches(0.15),
        y + Inches(0.07),
        w - Inches(0.2),
        Inches(0.3),
        [(h, dict(size=11, bold=True, color=ACCENT_AMBER))],
    )
y += Inches(0.4)
row_h = Inches(0.85)
for label, them, us in rows:
    add_rect(s, Inches(0.6), y, Inches(12.13), row_h, fill=LIGHT_BG, line=LINE)
    add_text(
        s,
        Inches(0.78),
        y + Inches(0.12),
        Inches(2.5),
        row_h - Inches(0.2),
        [(label, dict(size=11.5, bold=True, color=NAVY))],
    )
    add_text(
        s,
        Inches(3.45),
        y + Inches(0.12),
        Inches(4.5),
        row_h - Inches(0.2),
        [(them, dict(size=10.5, color=GREY))],
    )
    add_text(
        s,
        Inches(8.17),
        y + Inches(0.12),
        Inches(4.5),
        row_h - Inches(0.2),
        [(us, dict(size=10.5, color=INK, bold=True))],
    )
    y += row_h + Inches(0.04)
add_footer(s, 13, SLIDES_PLANNED)

# 14. Common questions ----------------------------------------------------
s = prs.slides.add_slide(blank)
add_slide_title(
    s,
    "Common questions from security teams",
    "Five things we hear most often during evaluation.",
)
qa = [
    (
        "“Why not just rotate keys per agent in our existing IAM?”",
        "You can. But Okta and Entra are session-oriented; an agent isn't a session — it's a long-lived process making thousands of decisions. Building the policy engine and the forensic chain on top of human IAM is a full quarter of platform work.",
    ),
    (
        "“What's the latency on the hot path?”",
        "Gateway target is <50ms p99 per enforcement decision. Verification of attestations happens offline — runtime systems never block on AI Identity. Designed to scale to millions of agent calls per second.",
    ),
    (
        "“How does this work with our SIEM?”",
        "Audit export API. Push to Splunk, Sentinel, Datadog, or a custom sink. The signed log remains the source of truth; the SIEM is the analytics and alerting layer.",
    ),
    (
        "“What if AI Identity itself gets compromised?”",
        "The HMAC chain key and the KMS signing key are held in separate systems — application database and FIPS 140-2 Level 3 HSM. Forging history requires both. The KMS private key has never been seen by an AI Identity process; auditors can verify against Google Cloud KMS directly.",
    ),
    (
        "“What does this cost?”",
        "Free $0, Pro $79, Business $299, Business+ $599 — rate-card published at ai-identity.co/pricing. Enterprise is $1,500/mo multi-tenant or $3,000/mo dedicated VPC. Standard Motor specifically: see the founding customer offer on the next page.",
    ),
]
y = Inches(2.05)
for q, a in qa:
    add_rect(s, Inches(0.6), y, Inches(12.13), Inches(0.95), fill=LIGHT_BG, line=LINE)
    add_rect(s, Inches(0.6), y, Emu(60000), Inches(0.95), fill=ACCENT_BLUE)
    add_text(
        s,
        Inches(0.85),
        y + Inches(0.1),
        Inches(11.7),
        Inches(0.4),
        [(q, dict(size=12, bold=True, color=NAVY))],
    )
    add_text(
        s,
        Inches(0.85),
        y + Inches(0.45),
        Inches(11.7),
        Inches(0.5),
        [(a, dict(size=10.5, color=INK))],
    )
    y += Inches(1.0)
add_footer(s, 14, SLIDES_PLANNED)

# 15. Standard Motor founding customer offer ------------------------------
s = prs.slides.add_slide(blank)
add_slide_title(
    s,
    "Standard Motor founding customer offer",
    "A specific, time-bounded path forward — should it be useful.",
)
# Hero card with the offer
add_rect(s, Inches(0.6), Inches(2.05), Inches(12.13), Inches(1.4), fill=NAVY, line=None)
add_rect(s, Inches(0.6), Inches(2.05), Inches(0.18), Inches(1.4), fill=ACCENT_AMBER)
add_text(
    s,
    Inches(1.0),
    Inches(2.18),
    Inches(11.5),
    Inches(0.4),
    [("FOUNDING ENTERPRISE CUSTOMER", dict(size=11, bold=True, color=ACCENT_AMBER))],
)
add_text(
    s,
    Inches(1.0),
    Inches(2.55),
    Inches(11.5),
    Inches(0.6),
    [
        (
            "90 days of the Enterprise tier — free, no commitment.",
            dict(size=22, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF)),
        )
    ],
)
add_text(
    s,
    Inches(1.0),
    Inches(3.05),
    Inches(11.5),
    Inches(0.4),
    [
        (
            "If it's not a fit at day 90, walk away. No termination clause, no kickback.",
            dict(size=12, color=RGBColor(0xC9, 0xD4, 0xE6)),
        )
    ],
)

# Three columns: What's included / What you get as founding customer / What happens at day 90
col_w = Inches(3.94)
col_h = Inches(3.05)
col_y = Inches(3.65)
gap = Inches(0.16)
col1_x = Inches(0.6)
col2_x = col1_x + col_w + gap
col3_x = col2_x + col_w + gap

# Column 1: What's included (Enterprise tier)
add_rect(s, col1_x, col_y, col_w, col_h, fill=LIGHT_BG, line=LINE)
add_rect(s, col1_x, col_y, col_w, Inches(0.12), fill=ACCENT_PURPLE)
add_text(
    s,
    col1_x + Inches(0.2),
    col_y + Inches(0.22),
    col_w - Inches(0.3),
    Inches(0.4),
    [("Enterprise tier — included", dict(size=13, bold=True, color=NAVY))],
)
ent_features = [
    "Unlimited agents",
    "Unlimited requests",
    "Unlimited audit retention",
    "Full SSO + SAML / SCIM",
    "Compliance evidence export",
    "Team roles + agent assignments",
    "Human-in-the-loop review",
    "On-premise / VPC deployment",
    "Dedicated support + SLA",
]
for i, f in enumerate(ent_features):
    add_text(
        s,
        col1_x + Inches(0.22),
        col_y + Inches(0.7) + Inches(0.25) * i,
        col_w - Inches(0.4),
        Inches(0.25),
        [(f"·  {f}", dict(size=10, color=INK))],
    )

# Column 2: What founding customer means
add_rect(s, col2_x, col_y, col_w, col_h, fill=LIGHT_BG, line=LINE)
add_rect(s, col2_x, col_y, col_w, Inches(0.12), fill=ACCENT_GREEN)
add_text(
    s,
    col2_x + Inches(0.2),
    col_y + Inches(0.22),
    col_w - Inches(0.3),
    Inches(0.4),
    [("What founding-customer status gets you", dict(size=13, bold=True, color=NAVY))],
)
fc_perks = [
    "Direct founder line — Slack or email, same-day responsiveness.",
    "Real influence on the product roadmap during your evaluation.",
    "First three policies authored with you, not by you.",
    "Quarterly compliance evidence export through your auditor of choice.",
    "Optional named launch case study — full sign-off on every quote.",
    "Pricing locked at evaluation rates if you continue past day 90.",
]
for i, line in enumerate(fc_perks):
    add_text(
        s,
        col2_x + Inches(0.22),
        col_y + Inches(0.7) + Inches(0.37) * i,
        col_w - Inches(0.4),
        Inches(0.4),
        [(f"·  {line}", dict(size=10, color=INK))],
    )

# Column 3: What happens at day 90
add_rect(s, col3_x, col_y, col_w, col_h, fill=LIGHT_BG, line=LINE)
add_rect(s, col3_x, col_y, col_w, Inches(0.12), fill=ACCENT_BLUE)
add_text(
    s,
    col3_x + Inches(0.2),
    col_y + Inches(0.22),
    col_w - Inches(0.3),
    Inches(0.4),
    [("What happens at day 90", dict(size=13, bold=True, color=NAVY))],
)
day90 = [
    (
        "Continue on Enterprise — $1,500/mo",
        "Multi-tenant Enterprise, locked at this rate for 24 months. ~$18K/yr — well below typical discretionary thresholds.",
    ),
    (
        "Dedicated VPC — $3,000/mo",
        "Same Enterprise feature set on its own infrastructure if compliance posture requires it. Locked for 24 months.",
    ),
    (
        "Step down — $79 / $299 / $599",
        "Pro, Business, or Business+ tier if real usage doesn't justify Enterprise. Or walk away — no termination clause.",
    ),
]
yy = col_y + Inches(0.7)
for label, body in day90:
    add_text(
        s,
        col3_x + Inches(0.22),
        yy,
        col_w - Inches(0.4),
        Inches(0.3),
        [(label, dict(size=11, bold=True, color=ACCENT_BLUE))],
    )
    add_text(
        s,
        col3_x + Inches(0.22),
        yy + Inches(0.3),
        col_w - Inches(0.4),
        Inches(0.5),
        [(body, dict(size=9.5, color=INK))],
    )
    yy += Inches(0.85)

add_text(
    s,
    Inches(0.6),
    Inches(6.85),
    Inches(12.13),
    Inches(0.3),
    [
        (
            "Honest framing: this is the offer because Standard Motor is exactly the kind of regulated, agent-curious enterprise we built this for.",
            dict(size=10.5, color=GREY, bold=True),
        )
    ],
)
add_footer(s, 15, SLIDES_PLANNED)


# 16. Verify it yourself / contact ----------------------------------------
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, prs.slide_width, prs.slide_height, fill=NAVY)
add_text(
    s,
    Inches(1.0),
    Inches(0.6),
    Inches(11.3),
    Inches(0.4),
    [("FURTHER READING", dict(size=14, bold=True, color=ACCENT_AMBER))],
)
add_text(
    s,
    Inches(1.0),
    Inches(1.05),
    Inches(11.3),
    Inches(0.7),
    [
        (
            "Verify everything in this deck independently",
            dict(size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF)),
        )
    ],
)
verify_items = [
    (
        "Public JWKS endpoint",
        "ai-identity.co/.well-known/ai-identity-public-keys.json — every non-destroyed forensic signing key version, fetchable without an account.",
    ),
    (
        "Open-source verification CLI",
        "github.com/Levaj2000/AI-Identity — clone and run cli/ai_identity_verify.py against any attestation envelope. Reproduces the trust chain offline.",
    ),
    (
        "Trust model (auditor reference)",
        "ai-identity.co/docs/forensics/trust-model — exhaustively documents what an attestation does and does not prove.",
    ),
    (
        "Architecture overview",
        "ai-identity.co/docs/architecture — gateway pipeline, latency budgets, deployment topology.",
    ),
    (
        "Compliance export profiles",
        "SOC 2 Type II  ·  EU AI Act  ·  NIST AI RMF  ·  GDPR — schema and audit-evidence format published in advance.",
    ),
]
y = Inches(2.1)
for label, body in verify_items:
    add_rect(
        s, Inches(1.0), y, Inches(11.3), Inches(0.65), fill=RGBColor(0x14, 0x23, 0x42), line=None
    )
    add_rect(s, Inches(1.0), y, Emu(60000), Inches(0.65), fill=ACCENT_AMBER)
    add_text(
        s,
        Inches(1.18),
        y + Inches(0.06),
        Inches(3.6),
        Inches(0.3),
        [(label, dict(size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF)))],
    )
    add_text(
        s,
        Inches(1.18),
        y + Inches(0.34),
        Inches(11.0),
        Inches(0.3),
        [(body, dict(size=10, color=RGBColor(0xC9, 0xD4, 0xE6)))],
    )
    y += Inches(0.72)
add_rect(s, Inches(1.0), Inches(6.05), Inches(11.3), Pt(1), fill=ACCENT_AMBER)
add_text(
    s,
    Inches(1.0),
    Inches(6.18),
    Inches(11.3),
    Inches(0.4),
    [("CONTACT", dict(size=11, bold=True, color=ACCENT_AMBER))],
)
add_text(
    s,
    Inches(1.0),
    Inches(6.5),
    Inches(11.3),
    Inches(0.5),
    [
        (
            f"{FOUNDER_NAME}  ·  Founder & CEO  ·  {FOUNDER_EMAIL}  ·  {WEBSITE}",
            dict(size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF)),
        )
    ],
)
add_text(
    s,
    Inches(0.6),
    Inches(7.18),
    Inches(12.13),
    Inches(0.3),
    [
        (
            f"{DECK_FOOTER}  •  16 / {SLIDES_PLANNED}",
            dict(size=9, color=RGBColor(0x6F, 0x7D, 0x9C), align="right"),
        )
    ],
)

# Save -------------------------------------------------------------------
OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"Slides: {len(prs.slides)}")

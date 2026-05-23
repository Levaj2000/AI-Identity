"""Build the Range Ventures investor deck.

AI Identity Purple template — first-meeting deck for Hayfa Aboukier (Range
Ventures), sent in response to her 2026-05-12 request for materials.

Run: python3 build_range_ventures_deck.py
Output: ../range-ventures-deck-2026-05-12.pptx
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── AI Identity Purple theme ─────────────────────────────────────────────

PURPLE_DEEP = RGBColor(0x4B, 0x2D, 0x7F)  # primary
PURPLE_MID = RGBColor(0x6B, 0x4F, 0xA0)  # secondary/accent
LAVENDER = RGBColor(0xF0, 0xE8, 0xF8)  # callout/KPI cards
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1A, 0x1F, 0x36)  # body text
GREY = RGBColor(0x6B, 0x72, 0x80)  # captions
LINE = RGBColor(0xE2, 0xE0, 0xEE)  # subtle borders

FONT = "Inter"
FONT_FALLBACK = "Calibri"

DECK_TITLE = "AI Identity"
DECK_SUBTITLE = "Trust Root for the Agent Economy"
DECK_FOOTER = "AI Identity  ·  Investor briefing  ·  Range Ventures"
TODAY = "May 12, 2026"
PREPARED_FOR = "Hayfa Aboukier · Range Ventures"
FOUNDER_EMAIL = "jeff@ai-identity.co"
FOUNDER_NAME = "Jeff Leva"
WEBSITE = "ai-identity.co"

OUT = Path(__file__).resolve().parent.parent / "range-ventures-deck-2026-05-12.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
EDGE_BAR_H = Emu(70000)  # ~6-8px purple top bar

# ── Helpers ──────────────────────────────────────────────────────────────


def set_text(tf, runs, *, anchor=MSO_ANCHOR.TOP):
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    for i, (text, opts) in enumerate(runs):
        if opts.get("align") == "inline" and i > 0:
            r = p.add_run()
        else:
            if i > 0:
                p = tf.add_paragraph()
            r = p.runs[0] if p.runs else p.add_run()
        if opts.get("align") in ("center", "left", "right"):
            p.alignment = {
                "center": PP_ALIGN.CENTER,
                "left": PP_ALIGN.LEFT,
                "right": PP_ALIGN.RIGHT,
            }[opts["align"]]
        r.text = text
        r.font.name = opts.get("font", FONT)
        r.font.size = Pt(opts.get("size", 14))
        r.font.bold = opts.get("bold", False)
        r.font.italic = opts.get("italic", False)
        r.font.color.rgb = opts.get("color", INK)
        if "space_after" in opts:
            p.space_after = Pt(opts["space_after"])


def add_rect(slide, x, y, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    return s


def add_text(slide, x, y, w, h, runs, *, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    set_text(box.text_frame, runs, anchor=anchor)
    return box


def add_circle_number(slide, x, y, diam, number):
    """Filled purple circle with white numeral — the brand's numbered list bullet."""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, diam, diam)
    c.shadow.inherit = False
    c.fill.solid()
    c.fill.fore_color.rgb = PURPLE_DEEP
    c.line.fill.background()
    tf = c.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    set_text(
        tf,
        [(str(number), dict(size=14, bold=True, color=WHITE, align="center"))],
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_callout(slide, x, y, w, h, label, body):
    """Lavender callout box ('LEARN' / 'WHY THIS MATTERS' pattern)."""
    add_rect(slide, x, y, w, h, fill=LAVENDER, line=None)
    add_text(
        slide,
        x + Inches(0.2),
        y + Inches(0.1),
        w - Inches(0.4),
        Inches(0.3),
        [(label, dict(size=10, bold=True, color=PURPLE_DEEP))],
    )
    add_text(
        slide,
        x + Inches(0.2),
        y + Inches(0.4),
        w - Inches(0.4),
        h - Inches(0.5),
        [(body, dict(size=12, color=INK))],
    )


def add_kpi_card(slide, x, y, w, h, *, number, label, sublabel=None):
    add_rect(slide, x, y, w, h, fill=LAVENDER, line=None)
    add_text(
        slide,
        x + Inches(0.2),
        y + Inches(0.2),
        w - Inches(0.4),
        Inches(0.9),
        [(number, dict(size=32, bold=True, color=PURPLE_DEEP))],
    )
    add_text(
        slide,
        x + Inches(0.2),
        y + Inches(1.15),
        w - Inches(0.4),
        Inches(0.45),
        [(label, dict(size=12, bold=True, color=INK))],
    )
    if sublabel:
        add_text(
            slide,
            x + Inches(0.2),
            y + Inches(1.55),
            w - Inches(0.4),
            Inches(0.5),
            [(sublabel, dict(size=10, color=GREY, italic=True))],
        )


def chrome(slide, *, eyebrow, title, page_n, total_pages):
    """Universal slide chrome: top purple bar, eyebrow, title, footer, page n/N."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, EDGE_BAR_H, fill=PURPLE_DEEP)
    add_text(
        slide,
        Inches(0.5),
        Inches(0.35),
        Inches(10),
        Inches(0.3),
        [(eyebrow.upper(), dict(size=10, bold=True, color=PURPLE_MID))],
    )
    add_text(
        slide,
        Inches(0.5),
        Inches(0.7),
        Inches(12.3),
        Inches(0.8),
        [(title, dict(size=28, bold=True, color=PURPLE_DEEP))],
    )
    # Footer
    add_text(
        slide,
        Inches(0.5),
        SLIDE_H - Inches(0.4),
        Inches(9),
        Inches(0.3),
        [(DECK_FOOTER, dict(size=9, color=GREY))],
    )
    add_text(
        slide,
        SLIDE_W - Inches(1.5),
        SLIDE_H - Inches(0.4),
        Inches(1),
        Inches(0.3),
        [(f"{page_n} / {total_pages}", dict(size=9, color=GREY, align="right"))],
    )


# ── Slide builders ───────────────────────────────────────────────────────


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    # Full purple background
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=PURPLE_DEEP)
    # Left accent bar (lighter purple)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=PURPLE_MID)
    # Title block
    add_text(
        s,
        Inches(0.8),
        Inches(2.2),
        Inches(11),
        Inches(1.4),
        [(DECK_TITLE, dict(size=72, bold=True, color=WHITE))],
    )
    add_text(
        s,
        Inches(0.8),
        Inches(3.6),
        Inches(11),
        Inches(0.8),
        [(DECK_SUBTITLE, dict(size=28, color=LAVENDER))],
    )
    # Thin lavender accent rule
    add_rect(s, Inches(0.8), Inches(4.5), Inches(2.5), Emu(20000), fill=LAVENDER)
    # Prepared-for block bottom-left
    add_text(
        s,
        Inches(0.8),
        Inches(6.3),
        Inches(8),
        Inches(0.4),
        [("PREPARED FOR", dict(size=10, bold=True, color=PURPLE_MID))],
    )
    add_text(
        s,
        Inches(0.8),
        Inches(6.6),
        Inches(8),
        Inches(0.4),
        [(PREPARED_FOR, dict(size=14, color=WHITE))],
    )
    add_text(
        s,
        SLIDE_W - Inches(3.5),
        Inches(6.6),
        Inches(3),
        Inches(0.4),
        [(TODAY, dict(size=12, color=LAVENDER, align="right"))],
    )


def slide_problem(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(
        s,
        eyebrow="The problem",
        title="Every AI agent today runs on a shared key.",
        page_n=n,
        total_pages=total,
    )

    body = (
        "Enterprises are moving fast into agentic AI. Agents are about to take real "
        "actions on real systems — create records, move data, call APIs, spend money "
        "on the company's behalf."
    )
    add_text(
        s, Inches(0.5), Inches(1.9), Inches(7.5), Inches(1.5), [(body, dict(size=14, color=INK))]
    )

    # Three questions card
    add_text(
        s,
        Inches(0.5),
        Inches(3.5),
        Inches(7.5),
        Inches(0.4),
        [
            (
                "Three questions no enterprise can answer about an agent action today:",
                dict(size=13, bold=True, color=INK),
            )
        ],
    )

    questions = [
        ("Which agent", "took this action?"),
        ("Under whose authority", "was it authorized?"),
        ("Against which policy", "was it evaluated?"),
    ]
    for i, (bold_part, rest) in enumerate(questions):
        y = Inches(4.1 + i * 0.55)
        add_circle_number(s, Inches(0.5), y, Inches(0.4), i + 1)
        add_text(
            s,
            Inches(1.05),
            y + Inches(0.02),
            Inches(7),
            Inches(0.4),
            [
                (bold_part, dict(size=14, bold=True, color=PURPLE_DEEP, align="inline")),
                (f" {rest}", dict(size=14, color=INK, align="inline")),
            ],
        )

    # Right-side callout
    add_callout(
        s,
        Inches(8.6),
        Inches(2.0),
        Inches(4.3),
        Inches(3.4),
        "WHY THIS MATTERS",
        "The identity, policy, and audit primitives we built for humans and services "
        "over the last 20 years simply don't exist for agents. Every shared service-account "
        "key is a compliance bomb waiting to go off when EU AI Act enforcement begins or "
        "the first agent-caused incident lands in a SOC 2 audit.",
    )


def slide_why_now(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(
        s,
        eyebrow="Why now",
        title="Two forces collide — and the window is open.",
        page_n=n,
        total_pages=total,
    )

    # Two columns
    cards = [
        (
            "DEMAND",
            "Agentic AI in production",
            "Banks, healthcare systems, and government are running agents against real "
            "data this quarter, not next year. Every CISO is asking the same three "
            "questions and nobody has an answer.",
        ),
        (
            "REGULATION",
            "Compliance frameworks landing",
            "EU AI Act enforcement begins Aug 2026. NIST AI RMF and SOC 2 already "
            "name agent governance as a control area. Auditors will start asking "
            "specifically about agent identity, authorization, and traceability "
            "within 12 months.",
        ),
    ]
    for i, (eyebrow, title, body) in enumerate(cards):
        x = Inches(0.5 + i * 6.4)
        add_rect(s, x, Inches(2.0), Inches(6.0), Inches(3.5), fill=LAVENDER, line=None)
        add_rect(s, x, Inches(2.0), Inches(0.12), Inches(3.5), fill=PURPLE_DEEP)
        add_text(
            s,
            x + Inches(0.3),
            Inches(2.2),
            Inches(5.5),
            Inches(0.3),
            [(eyebrow, dict(size=10, bold=True, color=PURPLE_MID))],
        )
        add_text(
            s,
            x + Inches(0.3),
            Inches(2.55),
            Inches(5.5),
            Inches(0.5),
            [(title, dict(size=18, bold=True, color=PURPLE_DEEP))],
        )
        add_text(
            s,
            x + Inches(0.3),
            Inches(3.2),
            Inches(5.5),
            Inches(2.2),
            [(body, dict(size=13, color=INK))],
        )

    add_callout(
        s,
        Inches(0.5),
        Inches(5.85),
        Inches(12.3),
        Inches(1.0),
        "THE WINDOW",
        "Trust infrastructure for the agent economy gets built once. The companies "
        "that ship the primitives first become the layer everyone else builds on. "
        "We're 6–18 months ahead of where the buyer just started asking.",
    )


def slide_what_we_built(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(
        s,
        eyebrow="What we built",
        title="Three primitives, in production today.",
        page_n=n,
        total_pages=total,
    )

    items = [
        (
            "Cryptographic agent identity",
            "Every agent gets a signed, revocable, rotatable identity. Not a shared key. "
            "Issuance, rotation, revocation, TTL — all first-class.",
        ),
        (
            "Immutable, hash-chained audit log",
            "Each decision is cryptographically linked to the one before it at write time. "
            "Same primitive git uses for commit history — tamper-evident by construction. "
            "An auditor can replay the chain a year later and prove nothing was edited.",
        ),
        (
            "Context-aware policy enforcement",
            "Every agent decision is evaluated against the rules you set, at decision time. "
            "ABAC on agent metadata. Per-agent, per-decision — not per-service-account. "
            "The audit trail is the proof you enforced it.",
        ),
    ]
    for i, (title, body) in enumerate(items):
        y = Inches(1.95 + i * 1.55)
        add_circle_number(s, Inches(0.5), y, Inches(0.5), i + 1)
        add_text(
            s,
            Inches(1.2),
            y - Inches(0.02),
            Inches(11.5),
            Inches(0.5),
            [(title, dict(size=18, bold=True, color=PURPLE_DEEP))],
        )
        add_text(
            s,
            Inches(1.2),
            y + Inches(0.5),
            Inches(11.5),
            Inches(1.0),
            [(body, dict(size=13, color=INK))],
        )


def slide_three_horizons(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(
        s,
        eyebrow="Three horizons",
        title="From audit primitives to the trust root.",
        page_n=n,
        total_pages=total,
    )

    horizons = [
        (
            "H1",
            "Governance probes",
            "LIVE",
            "AI forensics, compliance packs, audit primitives. Live in production. "
            "Probes indexed for SEO. Pipeline in motion.",
        ),
        (
            "H2",
            "Mandate Service",
            "DEPLOYED INTERNALLY",
            "Cryptographically signed permission grants — the bridge from agent identity "
            "to commerce. Production-hardened today; gated on first design partner for "
            "public exposure. Crypto-agile (ECDSA today, ML-DSA-87 hybrid slot ready).",
        ),
        (
            "H3",
            "Trust root for the agent economy",
            "OPTIONALITY",
            "The signing, verification, and policy primitives that the rest of the agent "
            "economy builds on — the way Stripe became the payments layer or Okta became "
            "the identity layer for humans.",
        ),
    ]
    for i, (label, title, status, body) in enumerate(horizons):
        x = Inches(0.5 + i * 4.3)
        add_rect(s, x, Inches(1.9), Inches(4.0), Inches(4.5), fill=LAVENDER, line=None)
        add_rect(s, x, Inches(1.9), Inches(4.0), Inches(0.5), fill=PURPLE_DEEP)
        add_text(
            s,
            x + Inches(0.25),
            Inches(1.97),
            Inches(3.5),
            Inches(0.4),
            [(label, dict(size=14, bold=True, color=WHITE))],
        )
        add_text(
            s,
            x + Inches(0.25),
            Inches(2.55),
            Inches(3.5),
            Inches(0.6),
            [(title, dict(size=17, bold=True, color=PURPLE_DEEP))],
        )
        add_text(
            s,
            x + Inches(0.25),
            Inches(3.2),
            Inches(3.5),
            Inches(0.3),
            [(status, dict(size=9, bold=True, color=PURPLE_MID))],
        )
        add_text(
            s,
            x + Inches(0.25),
            Inches(3.55),
            Inches(3.5),
            Inches(2.7),
            [(body, dict(size=12, color=INK))],
        )


def slide_live_today(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(
        s,
        eyebrow="Live today",
        title="Production, hardened, dogfooded.",
        page_n=n,
        total_pages=total,
    )

    kpis = [
        ("Prod", "Running on GKE", "Binary Auth, Cloud Armor, RLS"),
        ("3", "Probes indexed", "/forensics, /compliance-pack, /pqc-readiness"),
        ("~20 mo", "Runway stacked", "Google for Startups + Atlas accelerator credits"),
        ("0 ms", "Added latency", "Agents verify offline; signing is off the hot path"),
    ]
    card_w = Inches(3.0)
    for i, (n_, lab, sub) in enumerate(kpis):
        x = Inches(0.5 + i * 3.2)
        add_kpi_card(s, x, Inches(2.0), card_w, Inches(2.2), number=n_, label=lab, sublabel=sub)

    # Crypto-agility callout below
    add_callout(
        s,
        Inches(0.5),
        Inches(4.6),
        Inches(12.3),
        Inches(2.0),
        "CRYPTO-AGILE BY DESIGN",
        "Every signature carries an algorithm identifier and a second-signature slot. "
        "We sign with ECDSA-P256 today and have the ML-DSA-87 post-quantum slot reserved. "
        "When NIST clocks tick over, hybrid signing rolls in without invalidating anything "
        "we've issued. Mandates issued today will still verify in 2030 — and in the "
        "post-quantum era, with no re-signing event.",
    )


def slide_dogfood(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(
        s,
        eyebrow="Dogfood proof",
        title="Our own agent fixed our own bug — today.",
        page_n=n,
        total_pages=total,
    )

    # Left column — the story
    body = (
        "Ada is our internal senior-engineer agent, built on Google ADK and deployed "
        "on Cloud Run. Every tool call she makes is authenticated with her own AI "
        "Identity key and emitted into our own audit log."
    )
    add_text(
        s, Inches(0.5), Inches(1.95), Inches(7.5), Inches(1.4), [(body, dict(size=14, color=INK))]
    )

    # Timeline
    timeline = [
        (
            "This morning",
            "Ada flagged a missing try/finally around a compliance-load-bearing trigger toggle.",
        ),
        ("Within the hour", "Ada drafted the fix; we verified her diff against the code on main."),
        (
            "Shipped today",
            "PR #263 merged, audit log shows every read and reasoning step Ada took.",
        ),
    ]
    for i, (when, what) in enumerate(timeline):
        y = Inches(3.5 + i * 0.85)
        add_circle_number(s, Inches(0.5), y, Inches(0.4), i + 1)
        add_text(
            s,
            Inches(1.05),
            y - Inches(0.02),
            Inches(7),
            Inches(0.4),
            [
                (when, dict(size=13, bold=True, color=PURPLE_DEEP, align="inline")),
                (f" — {what}", dict(size=13, color=INK, align="inline")),
            ],
        )

    # Right callout
    add_callout(
        s,
        Inches(8.6),
        Inches(2.0),
        Inches(4.3),
        Inches(4.7),
        "WHY ENGINEERS LEAN IN ON THIS SLIDE",
        'The first thing every technical buyer asks: "does the product actually work?" '
        "Our answer is the audit trail of our own agent doing real engineering work, "
        "signed by our own keys API, on our own production code. Closed loop. "
        "Most companies pitching agent infra can't show you this. We can show you "
        "today's commit.",
    )


def slide_traction(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(
        s,
        eyebrow="Traction",
        title="Moving from probes to design partners.",
        page_n=n,
        total_pages=total,
    )

    items = [
        (
            "Probes live and indexed",
            "/forensics, /industries/finance/compliance-pack, /pqc-readiness — "
            "all submitted to Search Console, all driving inbound signal.",
        ),
        (
            "Design-partner pipeline",
            "Active outreach in LinkedIn Sales Navigator. Refreshed one-pager going out "
            "this week. Targeting regulated industries first: finance, healthcare, gov.",
        ),
        (
            "Partner infrastructure on the table",
            "Google for Startups Cloud Program (formal acceptance 2026-04-24), "
            "MongoDB Atlas accelerator credits, Cisco partnership conversation in motion.",
        ),
        (
            "Capital efficiency",
            "Live in production on stacked credits. ~20 months of runway already in the "
            "bag before this raise. We're raising to compound speed, not to keep the lights on.",
        ),
    ]
    for i, (title, body) in enumerate(items):
        y = Inches(1.95 + i * 1.2)
        add_circle_number(s, Inches(0.5), y, Inches(0.4), i + 1)
        add_text(
            s,
            Inches(1.05),
            y - Inches(0.02),
            Inches(11.5),
            Inches(0.4),
            [(title, dict(size=14, bold=True, color=PURPLE_DEEP))],
        )
        add_text(
            s,
            Inches(1.05),
            y + Inches(0.4),
            Inches(11.5),
            Inches(0.7),
            [(body, dict(size=12, color=INK))],
        )


def slide_gtm(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(
        s,
        eyebrow="Go-to-market",
        title="Land via design partners. Expand via compliance.",
        page_n=n,
        total_pages=total,
    )

    items = [
        (
            "LAND",
            "Design-partner program",
            "One agent workflow, 30 days free, weekly 20-minute sync with the founder. "
            "Success criteria defined together in week 1. If it doesn't fit at day 30, "
            "we part as friends.",
        ),
        (
            "EXPAND",
            "Compliance artifact pull",
            "Once the design partner sees SOC 2 / EU AI Act / NIST AI RMF export profiles "
            'show up in their audit prep, AI Identity goes from "pilot" to "line item." '
            "The product becomes their compliance receipt.",
        ),
        (
            "AMPLIFY",
            "Reference + content loop",
            "Each successful partner generates a reference conversation with one other "
            "buyer. Content marketing (PQC whitepaper, compliance packs, probes) "
            "compounds the inbound.",
        ),
    ]
    for i, (label, title, body) in enumerate(items):
        x = Inches(0.5 + i * 4.3)
        add_rect(s, x, Inches(1.95), Inches(4.0), Inches(4.5), fill=LAVENDER, line=None)
        add_rect(s, x, Inches(1.95), Inches(4.0), Inches(0.5), fill=PURPLE_DEEP)
        add_text(
            s,
            x + Inches(0.25),
            Inches(2.02),
            Inches(3.5),
            Inches(0.4),
            [(label, dict(size=14, bold=True, color=WHITE))],
        )
        add_text(
            s,
            x + Inches(0.25),
            Inches(2.6),
            Inches(3.5),
            Inches(0.5),
            [(title, dict(size=16, bold=True, color=PURPLE_DEEP))],
        )
        add_text(
            s,
            x + Inches(0.25),
            Inches(3.3),
            Inches(3.5),
            Inches(3.0),
            [(body, dict(size=12, color=INK))],
        )


def slide_competition(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(
        s,
        eyebrow="Competition",
        title="Workflow vs. infrastructure — and a decade of production muscle.",
        page_n=n,
        total_pages=total,
    )

    # Top — competitive frame
    add_text(
        s,
        Inches(0.5),
        Inches(1.95),
        Inches(12.3),
        Inches(0.4),
        [
            (
                "Opal, Valence, Holistic AI, Cognition — all workflow products. "
                "AI Identity is the layer underneath whatever workflow you pick.",
                dict(size=14, color=INK),
            )
        ],
    )

    # Two-column: who we are NOT vs. who we ARE
    cards = [
        (
            "WHAT WE'RE NOT",
            "A workflow tool",
            "Not another agent builder. Not another orchestration platform. Not a "
            'vibes-based "AI safety" wrapper. Those are crowded categories with '
            "thin moats.",
        ),
        (
            "WHAT WE ARE",
            "Cryptographic infrastructure",
            "The signing, verification, and audit primitives below the workflow. "
            "Boring on purpose. Hard to replace once installed. Sticky like databases — "
            "not sticky like dashboards.",
        ),
    ]
    for i, (eyebrow, title, body) in enumerate(cards):
        x = Inches(0.5 + i * 6.4)
        add_rect(s, x, Inches(2.7), Inches(6.0), Inches(2.5), fill=LAVENDER, line=None)
        add_rect(s, x, Inches(2.7), Inches(0.12), Inches(2.5), fill=PURPLE_DEEP)
        add_text(
            s,
            x + Inches(0.3),
            Inches(2.85),
            Inches(5.5),
            Inches(0.3),
            [(eyebrow, dict(size=10, bold=True, color=PURPLE_MID))],
        )
        add_text(
            s,
            x + Inches(0.3),
            Inches(3.2),
            Inches(5.5),
            Inches(0.5),
            [(title, dict(size=18, bold=True, color=PURPLE_DEEP))],
        )
        add_text(
            s,
            x + Inches(0.3),
            Inches(3.8),
            Inches(5.5),
            Inches(1.3),
            [(body, dict(size=12, color=INK))],
        )

    # Why us callout
    add_callout(
        s,
        Inches(0.5),
        Inches(5.45),
        Inches(12.3),
        Inches(1.4),
        "WHY US",
        "12+ years operating production systems where ambiguity costs you — "
        "executive escalations at Sprint, a D- → A- operational turnaround at "
        "Google, an SRE seat on a banking platform with $50B+ AUM at FIS. "
        "Compliance and audit-trail instinct that doesn't show up in agent demos "
        "but absolutely shows up at audit time.",
    )


def slide_team(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(
        s,
        eyebrow="Team",
        title="Founder-led. Technical. Shipping daily.",
        page_n=n,
        total_pages=total,
    )

    # Founder name + role + Columbia chip on the right
    add_text(
        s,
        Inches(0.5),
        Inches(1.95),
        Inches(8),
        Inches(0.5),
        [("Jeff Leva — Founder & CEO", dict(size=22, bold=True, color=PURPLE_DEEP))],
    )

    # Columbia chip (top-right of slide body area)
    add_rect(s, Inches(9.4), Inches(2.0), Inches(3.4), Inches(0.55), fill=LAVENDER, line=None)
    add_rect(s, Inches(9.4), Inches(2.0), Inches(0.08), Inches(0.55), fill=PURPLE_DEEP)
    add_text(
        s,
        Inches(9.55),
        Inches(2.05),
        Inches(3.2),
        Inches(0.25),
        [("INCOMING · JUNE 2026", dict(size=8, bold=True, color=PURPLE_MID))],
    )
    add_text(
        s,
        Inches(9.55),
        Inches(2.27),
        Inches(3.2),
        Inches(0.3),
        [("Columbia Business School CEO Program", dict(size=11, bold=True, color=PURPLE_DEEP))],
    )

    # Narrative paragraph
    body = (
        "12+ years operating production systems where reliability is the product. "
        "The through-line: roles where ambiguity costs you — executive escalations, "
        "turnarounds, and SRE seats in banking. AI Identity is the platform I kept "
        "wishing existed every time I found an agent acting on shared credentials "
        "with no audit trail."
    )
    add_text(
        s, Inches(0.5), Inches(2.7), Inches(12.3), Inches(1.2), [(body, dict(size=13, color=INK))]
    )

    # Career arc — four compact pills in a row
    arc = [
        (
            "FIS · SRE",
            "2022 – Present",
            "Cloud banking, $50B+ AUM, 99.9% uptime. SLO governance; 30% toil reduction via Python automation; 20% YoY drop in recurring criticals.",
        ),
        (
            "Google · Lead",
            "2019 – 2021",
            "Program lead + executive escalations. BBB rating D- → A- for 40+ consecutive weeks. 6.6x process-latency reduction.",
        ),
        (
            "British Telecom",
            "2018 – 2019",
            "Technical lead on Ernst & Young's video-conferencing infrastructure. 100% uptime across H.323/SIP.",
        ),
        (
            "Sprint · Escalations",
            "2014 – 2018",
            "Technical Analyst for executive escalations reporting to SVP. Root-cause on enterprise wireless, software, and network.",
        ),
    ]
    card_w = Inches(3.05)
    for i, (title, dates, body) in enumerate(arc):
        x = Inches(0.5 + i * 3.2)
        add_rect(s, x, Inches(4.0), card_w, Inches(2.4), fill=LAVENDER, line=None)
        add_rect(s, x, Inches(4.0), card_w, Inches(0.35), fill=PURPLE_DEEP)
        add_text(
            s,
            x + Inches(0.18),
            Inches(4.04),
            card_w - Inches(0.36),
            Inches(0.3),
            [(title, dict(size=11, bold=True, color=WHITE))],
        )
        add_text(
            s,
            x + Inches(0.18),
            Inches(4.45),
            card_w - Inches(0.36),
            Inches(0.3),
            [(dates, dict(size=9, color=PURPLE_MID, bold=True))],
        )
        add_text(
            s,
            x + Inches(0.18),
            Inches(4.75),
            card_w - Inches(0.36),
            Inches(1.6),
            [(body, dict(size=10, color=INK))],
        )

    # Capital-efficiency one-liner (no callout — keep it to one slim line so it
    # doesn't crowd the career arc or run into the footer at 7.1")
    add_text(
        s,
        Inches(0.5),
        Inches(6.55),
        Inches(12.3),
        Inches(0.4),
        [
            (
                "Built solo on stacked Google for Startups + MongoDB Atlas credits — "
                "~20 months of runway in the bag before this raise.",
                dict(size=11, italic=True, color=PURPLE_DEEP),
            )
        ],
    )


def slide_ask(prs, n, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    chrome(
        s,
        eyebrow="The ask",
        title="$1.25M SAFE → 18 months → 1 senior eng + GTM engine.",
        page_n=n,
        total_pages=total,
    )

    items = [
        ("$1.25M SAFE", "$8M post-money cap. Range as the lead."),
        (
            "18 months runway",
            "On top of the existing ~20 months of credits — true 38-month horizon.",
        ),
        (
            "Use of funds",
            "One senior engineer (mandate service productization + design-partner onboarding). "
            "Top-of-funnel ads & content — paid LinkedIn, sponsored newsletters, regulated-industry "
            "conferences. Design-partner conversion engine. SOC 2 Type II prep.",
        ),
        (
            "Milestones",
            "First 3 design partners closed by Q4 2026. Mandate Service publicly exposed with "
            "a named first user. Seed round opened mid-2027 with revenue + reference customers.",
        ),
    ]
    for i, (title, body) in enumerate(items):
        y = Inches(1.95 + i * 1.15)
        add_circle_number(s, Inches(0.5), y, Inches(0.4), i + 1)
        add_text(
            s,
            Inches(1.05),
            y - Inches(0.02),
            Inches(11.5),
            Inches(0.4),
            [(title, dict(size=15, bold=True, color=PURPLE_DEEP))],
        )
        add_text(
            s,
            Inches(1.05),
            y + Inches(0.4),
            Inches(11.5),
            Inches(0.65),
            [(body, dict(size=12, color=INK))],
        )


def slide_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=PURPLE_DEEP)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), SLIDE_H, fill=PURPLE_MID)

    add_text(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(11.5),
        Inches(2.0),
        [("The agent economy needs a trust root.", dict(size=40, bold=True, color=WHITE))],
    )
    add_text(
        s,
        Inches(0.8),
        Inches(4.0),
        Inches(11),
        Inches(0.7),
        [("We're building it.", dict(size=28, italic=True, color=LAVENDER))],
    )
    add_rect(s, Inches(0.8), Inches(4.85), Inches(2.5), Emu(20000), fill=LAVENDER)

    # Contact
    add_text(
        s,
        Inches(0.8),
        Inches(5.5),
        Inches(8),
        Inches(0.4),
        [("CONTACT", dict(size=10, bold=True, color=PURPLE_MID))],
    )
    add_text(
        s,
        Inches(0.8),
        Inches(5.85),
        Inches(8),
        Inches(0.4),
        [(FOUNDER_NAME, dict(size=16, bold=True, color=WHITE))],
    )
    add_text(
        s,
        Inches(0.8),
        Inches(6.2),
        Inches(8),
        Inches(0.4),
        [(FOUNDER_EMAIL, dict(size=14, color=LAVENDER))],
    )
    add_text(
        s,
        Inches(0.8),
        Inches(6.55),
        Inches(8),
        Inches(0.4),
        [(WEBSITE, dict(size=14, color=LAVENDER))],
    )


# ── Main ─────────────────────────────────────────────────────────────────


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Cover + 11 content + closing = 13 total
    total = 13

    slide_cover(prs)
    slide_problem(prs, 2, total)
    slide_why_now(prs, 3, total)
    slide_what_we_built(prs, 4, total)
    slide_three_horizons(prs, 5, total)
    slide_live_today(prs, 6, total)
    slide_dogfood(prs, 7, total)
    slide_traction(prs, 8, total)
    slide_gtm(prs, 9, total)
    slide_competition(prs, 10, total)
    slide_team(prs, 11, total)
    slide_ask(prs, 12, total)
    slide_closing(prs)

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
